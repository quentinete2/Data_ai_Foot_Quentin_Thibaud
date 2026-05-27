"""
Génère presentation_wc2026.pptx — Projet BigData B3
Auteurs : Thibaud & Quentin
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

BASE = os.path.dirname(os.path.abspath(__file__))

# ── Palette ──────────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1A, 0x3C, 0x6E)   # bleu FIFA profond
BLUE_MED   = RGBColor(0x1E, 0x5F, 0xAC)   # bleu moyen
GOLD       = RGBColor(0xD4, 0xAF, 0x37)   # or
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xFA)
DARK_GREY  = RGBColor(0x2D, 0x2D, 0x2D)
GREEN      = RGBColor(0x27, 0xAE, 0x60)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # totalement vide

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, x, y, w, h,
             size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def title_bar(slide, title, subtitle=None):
    """Barre de titre bleue en haut."""
    add_rect(slide, 0, 0, W, Inches(1.5), BLUE_DARK)
    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.9),
             size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(1.0), Inches(12), Inches(0.4),
                 size=16, color=GOLD, align=PP_ALIGN.LEFT)
    # trait doré sous la barre
    add_rect(slide, 0, Inches(1.5), W, Inches(0.04), GOLD)

def bg_light(slide):
    add_rect(slide, 0, 0, W, H, LIGHT_GREY)

def bullet_block(slide, items, x, y, w, h, size=20, color=DARK_GREY, indent="  ●  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = indent + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def add_table(slide, headers, rows, x, y, w, h):
    cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, cols, x, y, w, h).table
    # header
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_DARK
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = hdr
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.size = Pt(14)
    # rows
    for ri, row in enumerate(rows):
        fill = LIGHT_GREY if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = val
            run.font.size = Pt(13)
            run.font.color.rgb = DARK_GREY

# ── Slide 1 — Titre ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, BLUE_DARK)
# dégradé simulé avec une bande dorée diagonale
add_rect(s, 0, Inches(5.5), W, Inches(2.0), RGBColor(0x0F, 0x2A, 0x55))
add_rect(s, 0, Inches(7.3), W, Inches(0.2), GOLD)

add_text(s, "Prédiction FIFA",
         Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Coupe du Monde 2026",
         Inches(0.8), Inches(2.2), Inches(11), Inches(1.2),
         size=52, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "Analyse prédictive sur 92 ans de données (1930 – 2022)",
         Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
         size=22, bold=False, color=RGBColor(0xAD, 0xC6, 0xE8), align=PP_ALIGN.CENTER)
add_text(s, "Thibaud & Quentin  ·  Cours BigData B3  ·  2026",
         Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
         size=18, color=RGBColor(0xAD, 0xC6, 0xE8), align=PP_ALIGN.CENTER)

# ── Slide 2 — Question Prédictive ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Question Prédictive", "Notre objectif")

add_text(s, "Qui va gagner un match de Coupe du Monde ?",
         Inches(0.6), Inches(1.8), Inches(12), Inches(0.8),
         size=30, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

# 3 boîtes : home win / draw / away win
for i, (label, icon, col) in enumerate([
    ("Victoire\nÉquipe Dom.", "🏠", GREEN),
    ("Match Nul", "🤝", GOLD),
    ("Victoire\nÉquipe Ext.", "✈️", BLUE_MED),
]):
    bx = Inches(0.8 + i * 4.0)
    add_rect(s, bx, Inches(2.9), Inches(3.5), Inches(2.2), col)
    add_text(s, icon + "\n" + label,
             bx, Inches(2.9), Inches(3.5), Inches(2.2),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Tâche : Classification multiclasse  (3 sorties)",
         Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
         size=18, color=BLUE_DARK, align=PP_ALIGN.CENTER)

bullet_block(s, [
    "Modèle entraîné sur les matchs réels de 24 Coupes du Monde (1930–2022)",
    "84 équipes nationales actives — aucune rétro-ingénierie post-match",
    "Entrée : deux noms d'équipes  →  Sortie : résultat prédit + probabilités",
], Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.4), size=16)

# ── Slide 3 — Nos Données ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Nos Données", "Sources & périmètre")

# 2 cartes sources
for i, (title, body, col) in enumerate([
    ("matches.csv", "1 248 matchs\n1930 – 2022\n24 tournois\nColonnes clés : équipes,\nscores, date, tournament_id", BLUE_DARK),
    ("tournaments.csv", "30 tournois\nColonnes clés : année,\npays hôte, nb équipes\n(16 → 32 équipes)", BLUE_MED),
]):
    bx = Inches(0.6 + i * 6.4)
    add_rect(s, bx, Inches(1.7), Inches(5.8), Inches(3.2), col)
    add_text(s, title, bx + Inches(0.15), Inches(1.8), Inches(5.5), Inches(0.6),
             size=22, bold=True, color=GOLD)
    add_text(s, body, bx + Inches(0.15), Inches(2.5), Inches(5.5), Inches(2.0),
             size=17, color=WHITE)

add_text(s, "Nettoyage & normalisation",
         Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
         size=18, bold=True, color=BLUE_DARK)
bullet_block(s, [
    "Allemagne Ouest / Est  →  Allemagne unifiée",
    "Union Soviétique  →  Russie",
    "Équipes défuntes exclues : Yougoslavie, Tchécoslovaquie, Sarre…",
    "Distribution des résultats : 677 victoires dom. · 253 nuls · 318 victoires ext.",
], Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.8), size=16)

# ── Slide 4 — Pipeline ETL ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Pipeline ETL", "Du CSV brut au modèle prêt à l'emploi")

steps = [
    ("1", "Chargement", "matches.csv\n+ tournaments.csv"),
    ("2", "Jointure", "Merge sur\ntournament_id"),
    ("3", "Nettoyage", "Noms d'équipes\nunifiés"),
    ("4", "Cible y", "0 dom. / 1 nul\n/ 2 ext."),
    ("5", "Features", "9 variables\npar match"),
    ("6", "Split", "80 % train\n20 % test"),
    ("7", "Export", "model.pkl\n(joblib)"),
]
box_w = Inches(1.6)
box_h = Inches(1.5)
gap   = Inches(0.15)
start_x = Inches(0.3)
y_box   = Inches(2.2)

for i, (num, title, body) in enumerate(steps):
    bx = start_x + i * (box_w + gap)
    col = BLUE_DARK if i < 6 else GREEN
    add_rect(s, bx, y_box, box_w, box_h, col)
    add_text(s, num, bx, y_box + Inches(0.05), box_w, Inches(0.35),
             size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, title, bx, y_box + Inches(0.35), box_w, Inches(0.45),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, body, bx, y_box + Inches(0.8), box_w, Inches(0.65),
             size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    # flèche sauf dernier
    if i < len(steps) - 1:
        ax = bx + box_w + Inches(0.02)
        add_text(s, "→", ax, y_box + Inches(0.55), gap, Inches(0.4),
                 size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

add_text(s, "Pondération temporelle : demi-vie exponentielle de 14 ans (référence 2022)",
         Inches(0.5), Inches(4.1), Inches(12), Inches(0.45),
         size=17, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_text(s,
    "Les matchs récents comptent plus : un match de 2022 pèse 2× plus qu'un match de 2008.\n"
    "Cette pondération est appliquée aux moyennes de buts et aux taux de victoire de chaque équipe.",
    Inches(0.8), Inches(4.65), Inches(11.5), Inches(0.9),
    size=15, color=DARK_GREY, align=PP_ALIGN.LEFT)

# ── Slide 5 — Features ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Ingénierie des Features", "9 variables construites par match")

headers = ["#", "Feature", "Type", "Définition"]
rows = [
    ("1", "home_avg_goals",         "float", "Moy. buts/match (dom.) — pondérée"),
    ("2", "away_avg_goals",         "float", "Moy. buts/match (ext.) — pondérée"),
    ("3", "goal_diff",              "float", "home_avg_goals − away_avg_goals"),
    ("4", "home_total_matches",     "int",   "Nb de matchs joués (dom.) en carrière"),
    ("5", "away_total_matches",     "int",   "Nb de matchs joués (ext.) en carrière"),
    ("6", "home_weighted_win_rate", "float", "% victoires (dom.) pondéré, [0–1]"),
    ("7", "away_weighted_win_rate", "float", "% victoires (ext.) pondéré, [0–1]"),
    ("8", "year",                   "int",   "Année du tournoi (médiane : 1993)"),
    ("9", "count_teams",            "int",   "Nb d'équipes dans le tournoi (médiane : 16)"),
]
add_table(s, headers, rows,
          Inches(0.4), Inches(1.65), Inches(12.5), Inches(4.9))

add_text(s,
    "Les features 8 & 9 utilisent la valeur médiane historique lors de l'inférence (aucune donnée 2026 disponible).",
    Inches(0.5), Inches(6.75), Inches(12), Inches(0.5),
    size=14, italic=True, color=BLUE_MED, align=PP_ALIGN.LEFT)

# ── Slide 6 — Modèle ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Choix du Modèle", "RandomForestClassifier")

# Boîte centrale modèle
add_rect(s, Inches(0.5), Inches(1.7), Inches(5.5), Inches(3.6), BLUE_DARK)
add_text(s, "RandomForestClassifier", Inches(0.65), Inches(1.8), Inches(5.2), Inches(0.7),
         size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
params = [
    "n_estimators = 100 arbres",
    "max_depth    = 10",
    "random_state = 42",
    "Train : 998 matchs (80 %)",
    "Test  : 250 matchs (20 %)",
]
bullet_block(s, params, Inches(0.65), Inches(2.55), Inches(5.2), Inches(2.5),
             size=17, color=WHITE, indent="  ▸  ")

add_text(s, "Pourquoi Random Forest ?",
         Inches(6.3), Inches(1.7), Inches(6.7), Inches(0.5),
         size=20, bold=True, color=BLUE_DARK)
reasons = [
    "Robuste aux valeurs aberrantes",
    "Capture les interactions non-linéaires entre features",
    "Fournit des probabilités calibrées (barres dans l'UI)",
    "Aucune normalisation requise (features déjà propres)",
    "Faible risque de surapprentissage avec max_depth=10",
    "Résistant aux données déséquilibrées (home win majoritaire)",
]
bullet_block(s, reasons, Inches(6.3), Inches(2.3), Inches(6.7), Inches(3.0),
             size=16, color=DARK_GREY)

# ── Slide 7 — Résultats ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Résultats du Modèle", "Évaluation sur le jeu de test (250 matchs)")

# Gros accuracy
add_rect(s, Inches(0.5), Inches(1.65), Inches(3.2), Inches(2.4), BLUE_DARK)
add_text(s, "64.80 %", Inches(0.5), Inches(1.75), Inches(3.2), Inches(1.2),
         size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "Accuracy\n(test set)", Inches(0.5), Inches(2.9), Inches(3.2), Inches(0.8),
         size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Tableau par classe
headers2 = ["Classe", "Précision", "Rappel", "F1", "Support"]
rows2 = [
    ("🏠 Victoire dom.  (0)", "0.70", "0.83", "0.76", "143"),
    ("🤝 Nul             (1)", "0.38", "0.19", "0.25 ⚠️", " 47"),
    ("✈️ Victoire ext.  (2)", "0.61", "0.57", "0.59", " 60"),
    ("Moyenne macro",          "0.56", "0.53", "0.53", "250"),
]
add_table(s, headers2, rows2,
          Inches(4.1), Inches(1.65), Inches(8.8), Inches(2.8))

add_text(s, "Points clés",
         Inches(0.5), Inches(4.2), Inches(12.5), Inches(0.4),
         size=18, bold=True, color=BLUE_DARK)
bullet_block(s, [
    "Les victoires à domicile sont bien prédites (F1 = 0.76) — elles représentent 57 % des matchs",
    "Les nuls sont difficiles à détecter : rappel de 19 % — classe très sous-représentée (19 % du corpus)",
    "Les victoires extérieures sont correctement classifiées (F1 = 0.59)",
], Inches(0.7), Inches(4.7), Inches(12), Inches(1.6), size=16)

# ── Slide 8 — Screenshot Dashboard ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Application — Dashboard", "Score de forme & statistiques comparatives")

sc_path = os.path.join(BASE, "screenshot_dashboard.png")
if os.path.exists(sc_path):
    s.shapes.add_picture(sc_path, Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.6))
else:
    add_rect(s, Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.6), BLUE_MED)
    add_text(s, "[Screenshot — Dashboard]\nInsère ici une capture d'écran du frontend",
             Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
             size=20, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Graphique : top 15 équipes classées par score composite (taux victoires + buts + expérience)",
         Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
         size=13, italic=True, color=BLUE_MED, align=PP_ALIGN.CENTER)

# ── Slide 9 — Screenshot Prédiction ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Application — Prédiction France vs Brazil", "Résultat en temps réel avec probabilités")

sc_path2 = os.path.join(BASE, "screenshot_prediction.png")
if os.path.exists(sc_path2):
    s.shapes.add_picture(sc_path2, Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.6))
else:
    add_rect(s, Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.6), BLUE_MED)
    add_text(s, "[Screenshot — Prédiction France vs Brazil]\nInsère ici la capture d'écran du résultat",
             Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
             size=20, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Résultat prédit + probabilités pour chaque issue (victoire dom. / nul / victoire ext.) + stats historiques des deux équipes",
         Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
         size=13, italic=True, color=BLUE_MED, align=PP_ALIGN.CENTER)

# ── Slide 10 — Pistes d'amélioration & Conclusion ────────────────────────────
s = prs.slides.add_slide(BLANK)
bg_light(s)
title_bar(s, "Pistes d'Amélioration & Conclusion", "")

# 2 colonnes
add_text(s, "Améliorations possibles",
         Inches(0.5), Inches(1.65), Inches(6.2), Inches(0.45),
         size=20, bold=True, color=BLUE_DARK)
improvements = [
    "Rééquilibrer les classes (SMOTE ou poids) → meilleur rappel des nuls",
    "Ajouter le classement FIFA / ELO comme feature",
    "Intégrer des stats joueurs (top buteurs, solidité défensive)",
    "Réduire la fenêtre temporelle (7–10 ans) pour les nations émergentes",
    "Simuler un tournoi complet (bracket WC 2026, 48 équipes)",
    "Tester Gradient Boosting (XGBoost/LightGBM) pour comparer",
]
bullet_block(s, improvements, Inches(0.5), Inches(2.2), Inches(6.2), Inches(4.8),
             size=15, color=DARK_GREY)

# Colonne droite — conclusion
add_rect(s, Inches(7.0), Inches(1.65), Inches(5.9), Inches(5.5), BLUE_DARK)
add_text(s, "Ce que nous avons accompli",
         Inches(7.15), Inches(1.75), Inches(5.6), Inches(0.5),
         size=18, bold=True, color=GOLD)
ccl = [
    "Pipeline ETL complet sur 92 ans de données FIFA",
    "9 features avec pondération temporelle",
    "Modèle à 64.8 % d'accuracy (3 classes)",
    "API FastAPI + Dashboard React interactif",
    "Déploiement Docker (frontend + backend)",
]
bullet_block(s, ccl, Inches(7.15), Inches(2.35), Inches(5.6), Inches(3.5),
             size=15, color=WHITE)

add_text(s, "Merci !",
         Inches(7.15), Inches(5.9), Inches(5.6), Inches(0.6),
         size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, "Thibaud & Quentin",
         Inches(7.15), Inches(6.5), Inches(5.6), Inches(0.4),
         size=16, color=WHITE, align=PP_ALIGN.CENTER)

# ── Export ────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "presentation_wc2026.pptx")
prs.save(out)
print(f"Fichier généré : {out}")
